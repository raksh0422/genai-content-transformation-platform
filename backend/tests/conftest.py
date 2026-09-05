"""Test fixtures and shared configuration."""
from __future__ import annotations

import io
import uuid
from pathlib import Path
from typing import AsyncGenerator

import pytest
import pytest_asyncio
from httpx import AsyncClient, ASGITransport
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine

from app.config import Settings
from app.database import Base, get_db
from app.models import Document, DocumentChunk  # noqa: F401 – register models


# ---------------------------------------------------------------------------
# In-memory SQLite settings (no Postgres needed for tests)
# ---------------------------------------------------------------------------

TEST_DATABASE_URL = "sqlite+aiosqlite:///:memory:"


@pytest.fixture(scope="session")
def test_settings(tmp_path_factory) -> Settings:
    """Return Settings pointing at a temp upload dir and SQLite."""
    upload_dir = tmp_path_factory.mktemp("uploads")
    return Settings(
        database_url=TEST_DATABASE_URL,
        upload_dir=upload_dir,
        max_file_size_mb=10,
        chunk_size_tokens=128,
        chunk_overlap_tokens=16,
        app_env="test",
        app_log_level="DEBUG",
    )


@pytest_asyncio.fixture(scope="session")
async def db_engine(test_settings):
    """Create the async engine and all tables once per test session."""
    engine = create_async_engine(test_settings.database_url, echo=False)
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    yield engine
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.drop_all)
    await engine.dispose()


@pytest_asyncio.fixture
async def db_session(db_engine) -> AsyncGenerator[AsyncSession, None]:
    """Yield a transactional session that's rolled back after each test."""
    factory = async_sessionmaker(db_engine, expire_on_commit=False)
    async with factory() as session:
        yield session
        await session.rollback()


@pytest_asyncio.fixture
async def client(test_settings, db_engine) -> AsyncGenerator[AsyncClient, None]:
    """Return an async test client with DB and settings dependencies overridden."""
    from app.main import create_app
    from app.config import get_settings

    app = create_app()

    # Override get_settings
    app.dependency_overrides[get_settings] = lambda: test_settings

    # Override get_db
    factory = async_sessionmaker(db_engine, expire_on_commit=False)

    async def override_get_db():
        async with factory() as session:
            try:
                yield session
                await session.commit()
            except Exception:
                await session.rollback()
                raise

    app.dependency_overrides[get_db] = override_get_db

    # Patch app.database globals so the lifespan startup doesn't try to
    # connect to Postgres or re-create a conflicting engine.
    import app.database as db_module
    original_engine = db_module._engine
    original_factory = db_module._session_factory
    db_module._engine = db_engine
    db_module._session_factory = factory

    try:
        async with AsyncClient(
            transport=ASGITransport(app=app), base_url="http://test"
        ) as ac:
            yield ac
    finally:
        db_module._engine = original_engine
        db_module._session_factory = original_factory


# ---------------------------------------------------------------------------
# Sample file fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_txt_bytes() -> bytes:
    return (
        b"Introduction\n\n"
        b"This is the first paragraph of the document. "
        b"It contains several sentences to provide enough content for testing.\n\n"
        b"Section Two\n\n"
        b"The second paragraph discusses more topics in detail. "
        b"We include enough words to trigger chunking behaviour.\n\n"
        b"Conclusion\n\n"
        b"A final concluding paragraph to wrap up the document."
    )


@pytest.fixture
def sample_pdf_bytes() -> bytes:
    """Generate a minimal valid PDF in memory using PyMuPDF."""
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Introduction", fontsize=18)
    page.insert_text((72, 120), "This is a sample paragraph for testing purposes.", fontsize=12)
    page.insert_text((72, 200), "Second Section", fontsize=16)
    page.insert_text((72, 240), "More content to test multi-block extraction.", fontsize=12)
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


@pytest.fixture
def sample_docx_bytes() -> bytes:
    """Generate a minimal DOCX in memory using python-docx."""
    from docx import Document as DocxDocument
    buf = io.BytesIO()
    doc = DocxDocument()
    doc.add_heading("Test Heading", level=1)
    doc.add_paragraph(
        "This is the first paragraph with enough words to be treated as body text."
    )
    doc.add_heading("Second Section", level=2)
    doc.add_paragraph("Another paragraph for section two of the test document.")
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    """Generate a minimal PPTX in memory using python-pptx."""
    from pptx import Presentation

    buf = io.BytesIO()
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide One Title"
    slide.placeholders[1].text = "Content of slide one for testing extraction."
    prs.save(buf)
    return buf.getvalue()
