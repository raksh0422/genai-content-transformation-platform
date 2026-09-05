from __future__ import annotations
"""PPTX parser using python-pptx."""
import logging
from io import BytesIO

from pptx import Presentation
from pptx.util import Pt

from app.services.processing.cleaner import clean_text
from app.services.processing.models import ExtractedBlock, ExtractionResult
from app.services.processing.structure import classify_block_type

logger = logging.getLogger(__name__)

# Font size threshold (in points) above which text is considered a heading
_HEADING_PT_THRESHOLD = 24.0


def _is_title_placeholder(placeholder) -> bool:
    """Return True if the placeholder is a title-type placeholder."""
    # Placeholder types: 0=body, 1=center_title, 2=subtitle, 3=date, 4=footer,
    #                    5=slide_number, 13=title, 15=vertical_title
    try:
        return placeholder.placeholder_format.idx in (1, 13, 15)
    except AttributeError:
        return False


def parse_pptx(content: bytes) -> ExtractionResult:
    """
    Extract text from a PPTX byte stream, preserving slide numbers.

    Title placeholders are tagged as headings. Body text is classified
    heuristically.

    Args:
        content: Raw PPTX file bytes.

    Returns:
        ExtractionResult with per-slide blocks.
    """
    result = ExtractionResult()

    try:
        prs = Presentation(BytesIO(content))
    except Exception as exc:
        logger.error("Failed to open PPTX: %s", exc)
        raise ValueError(f"Cannot parse PPTX: {exc}") from exc

    result.page_count = len(prs.slides)

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            is_title = _is_title_placeholder(shape)

            for para in shape.text_frame.paragraphs:
                raw_text = para.text
                cleaned = clean_text(raw_text)
                if not cleaned:
                    continue

                # Determine block type from placeholder type or font size
                if is_title:
                    block_type = "heading"
                else:
                    # Check run font sizes
                    max_size = 0.0
                    for run in para.runs:
                        if run.font.size:
                            size_pt = run.font.size / 12700  # EMUs to points
                            max_size = max(max_size, size_pt)

                    if max_size >= _HEADING_PT_THRESHOLD:
                        block_type = "heading"
                    else:
                        block_type = classify_block_type(cleaned)

                result.blocks.append(
                    ExtractedBlock(
                        text=cleaned,
                        block_type=block_type,
                        slide_number=slide_num,
                    )
                )

    result.word_count = sum(len(b.text.split()) for b in result.blocks)
    logger.info(
        "PPTX extraction complete: %d slides, %d blocks, %d words",
        result.page_count,
        len(result.blocks),
        result.word_count,
    )
    return result
